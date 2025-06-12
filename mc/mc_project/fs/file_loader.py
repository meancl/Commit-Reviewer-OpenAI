from pathlib import Path

def get_files_content(file_arg: str) -> str:
    result = []
    paths = file_arg.strip().split()
 
    for idx, path_str in enumerate(paths, 1):
        path = Path(path_str)

        if not path.exists() or not path.is_file():
            result.append(f"[ file{idx}. {path.name} ]\n<파일을 찾을 수 없습니다>\n")
            continue

        suffix = path.suffix.lower()

        try:
            if suffix == ".pdf":
                from PyPDF2 import PdfReader
                reader = PdfReader(str(path))
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
                result.append(f"[ file{idx}. {path.name} ]\n{text}\n")

            elif suffix == ".xlsx":
                from openpyxl import load_workbook
                wb = load_workbook(filename=path, read_only=True)
                text = ""
                for sheet in wb.worksheets:
                    text += f"\n[Sheet: {sheet.title}]\n"
                    for row in sheet.iter_rows(values_only=True):
                        line = "\t".join(str(cell) if cell is not None else "" for cell in row)
                        text += line + "\n"
                result.append(f"[ file{idx}. {path.name} ]\n{text}\n")

            elif suffix == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(path))
                text = ""
                for i, slide in enumerate(prs.slides):
                    text += f"\n[Slide {i + 1}]\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                result.append(f"[ file{idx}. {path.name} ]\n{text}\n")

            else:
                with path.open(encoding="utf-8") as f:
                    content = f.read()
                result.append(f"[ file{idx}. {path.name} ]\n{content}\n")
            
        except Exception as e:
            result.append(f"[ file{idx}. {path.name} ]\n<읽기 실패: {e}>\n")

    return "\n".join(result)